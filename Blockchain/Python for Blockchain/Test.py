from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

slides_content = [
("Operating System & Linux", "Overview and Introduction to OS & Linux\nTech Style Theme"),
("What is an Operating System?", 
 "• System software managing hardware & software resources\n• Acts as an interface between user and hardware\n• Examples: Windows, macOS, Linux"),
("Functions of an Operating System",
 "• Process Management\n• Memory Management\n• File System Management\n• Device Management\n• Security & Protection"),
("What is Linux?",
 "• Open-source UNIX-like operating system\n• Created by Linus Torvalds in 1991\n• Free to modify, use and distribute\n• Popular in servers, DevOps, cybersecurity, embedded systems"),
("Linux Architecture",
 "• User Applications\n• Shell\n• Kernel\n• Hardware"),
("Linux Kernel",
 "• Core part of the OS\n• Manages CPU, memory, devices & system calls\n• Types: Monolithic, Microkernel"),
("Linux File System Structure",
 "• Root directory (/)\n• /bin, /etc, /usr, /home, /var\n• Hierarchical tree structure"),
("Linux Distributions",
 "• Ubuntu\n• Fedora\n• Debian\n• Arch Linux\n• CentOS / RHEL"),
("Linux Shell & Commands",
 "• Shell interprets commands and executes them\n• Examples:\n  - ls, cd, mkdir, rm\n  - cat, nano, cp, mv"),
("Advantages of Linux",
 "• Open-source and free\n• Highly secure and stable\n• Customizable\n• Supports development & servers"),
("Real World Applications",
 "• Web servers (Apache, Nginx)\n• Cloud (AWS, Azure)\n• Cybersecurity tools (Kali Linux)\n• Android based on Linux kernel"),
("Conclusion",
 "Linux is powerful, secure and widely used in modern computing\nUnderstanding Linux is essential for developers and engineers")
]

for title, content in slides_content:
    add_slide(title, content)

ppt_path = "/mnt/data/Operating_System_Linux_v2.pptx"
prs.save(ppt_path)

ppt_path
